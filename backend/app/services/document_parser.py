"""
Universal Document Intelligence Extractor Module
Accepts ANY document format (PDF, DOCX, DOC, XLSX, XLS, CSV, TSV, PPTX, PPT, HTML, XML, JSON, RTF, TXT, MD, LOG, YAML, EPUB, ODT, etc.)
and extracts clean, structured text for AI analysis.
"""
import io
import re
from typing import Optional


def extract_text_from_any_document(filename: str, content: bytes, content_type: Optional[str] = None) -> str:
    """
    Extracts readable text content from ANY uploaded document file.
    Supports native parsing for PDF, Word, Excel, PowerPoint, HTML, XML, CSV, JSON,
    and falls back to robust multi-encoding text decoding for any arbitrary file type.
    """
    if not content:
        return ""

    fn = (filename or "").lower()
    extracted = ""

    # 1. PDF Documents (.pdf)
    if fn.endswith(".pdf") or content_type == "application/pdf":
        try:
            import pypdf
            reader = pypdf.PdfReader(io.BytesIO(content))
            extracted = "\n".join([page.extract_text() or "" for page in reader.pages if page.extract_text()])
        except Exception as e:
            print(f"[DocumentParser] pypdf error on {filename}: {e}")

    # 2. Word Documents (.docx, .doc)
    elif fn.endswith((".docx", ".doc")):
        try:
            import docx
            doc = docx.Document(io.BytesIO(content))
            paragraphs = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                    if row_text:
                        paragraphs.append(row_text)
            extracted = "\n".join(paragraphs)
        except Exception:
            # Fallback: Extract text from Word XML structure directly
            try:
                raw_xml = content.decode("utf-8", errors="ignore")
                extracted = "\n".join(re.findall(r"<w:t[^>]*>(.*?)</w:t>", raw_xml))
            except Exception:
                pass

    # 3. Excel Spreadsheets & Data Files (.xlsx, .xls, .csv, .tsv)
    elif fn.endswith((".xlsx", ".xls", ".csv", ".tsv")):
        if fn.endswith((".csv", ".tsv")):
            try:
                import csv
                delimiter = "\t" if fn.endswith(".tsv") else ","
                text_io = io.StringIO(content.decode("utf-8", errors="ignore"))
                reader = csv.reader(text_io, delimiter=delimiter)
                rows = ["\t".join(row) for row in reader if any(row)]
                extracted = "\n".join(rows)
            except Exception:
                pass

        if not extracted.strip():
            try:
                import pandas as pd
                df = pd.read_excel(io.BytesIO(content))
                extracted = df.to_string()
            except Exception:
                try:
                    import openpyxl
                    wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
                    lines = []
                    for sheet in wb.worksheets:
                        for row in sheet.iter_rows(values_only=True):
                            row_vals = [str(v) for v in row if v is not None]
                            if row_vals:
                                lines.append("\t".join(row_vals))
                    extracted = "\n".join(lines)
                except Exception:
                    pass

    # 4. PowerPoint Presentations (.pptx, .ppt)
    elif fn.endswith((".pptx", ".ppt")):
        try:
            import pptx
            prs = pptx.Presentation(io.BytesIO(content))
            lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        lines.append(shape.text.strip())
            extracted = "\n".join(lines)
        except Exception:
            pass

    # 5. HTML / XML / Rich Text (.html, .htm, .xml, .xhtml)
    elif fn.endswith((".html", ".htm", ".xml", ".xhtml")):
        try:
            text = content.decode("utf-8", errors="ignore")
            clean = re.sub(r"<style[^>]*>.*?</style>", "", text, flags=re.DOTALL | re.IGNORECASE)
            clean = re.sub(r"<script[^>]*>.*?</script>", "", clean, flags=re.DOTALL | re.IGNORECASE)
            clean = re.sub(r"<[^>]+>", " ", clean)
            extracted = "\n".join([line.strip() for line in clean.splitlines() if line.strip()])
        except Exception:
            pass

    # 6. JSON Data Files (.json)
    elif fn.endswith(".json"):
        try:
            import json
            obj = json.loads(content.decode("utf-8", errors="ignore"))
            extracted = json.dumps(obj, indent=2)
        except Exception:
            pass

    # 7. Universal Fallback for ANY Document Type (TXT, MD, RTF, YAML, LOG, UNKNOWN BINARY/TEXT)
    if not extracted.strip():
        for enc in ("utf-8", "utf-16", "latin-1", "ascii"):
            try:
                decoded = content.decode(enc, errors="ignore")
                # Remove unprintable control codes but keep tabs, newlines, readable unicode
                cleaned = re.sub(r"[^\x09\x0A\x0D\x20-\x7E\x80-\xFF]", " ", decoded)
                lines = [line.strip() for line in cleaned.splitlines() if line.strip()]
                if len(lines) > 0:
                    extracted = "\n".join(lines)
                    break
            except Exception:
                continue

    return extracted.strip()
